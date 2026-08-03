#!/usr/bin/env python3
"""Generate the Cadence hackathon checkpoint deck (16:9 .pptx).

Palette and layout mirror the live app so the deck reads as one product.
No em-dashes anywhere in the copy by design.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Brand palette (from tailwind.config.ts / globals.css) --------------------
ARC_BLUE   = RGBColor(0x16, 0x52, 0xF0)
ARC_DARK   = RGBColor(0x0A, 0x0B, 0x0D)
ARC_CARD   = RGBColor(0x11, 0x13, 0x18)
ARC_BORDER = RGBColor(0x1E, 0x22, 0x30)
WHITE      = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_400   = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500   = RGBColor(0x6B, 0x72, 0x80)
GREEN      = RGBColor(0x4A, 0xDE, 0x80)
USDC_BLUE  = RGBColor(0x27, 0x75, 0xCA)

FONT       = "Arial"          # safe cross-platform stand-in for Inter
MONO       = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---- helpers ------------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = ARC_DARK
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def _set_font(run, size, color, bold=False, mono=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = MONO if mono else FONT


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, space_after=0):
    """runs: list of paragraphs; each paragraph is a list of (txt,size,color,bold,mono) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        for seg in para:
            txt, size, color, bold, mono = (seg + (False, False))[:5]
            r = p.add_run(); r.text = txt
            _set_font(r, size, color, bold, mono)
    return tb


def card(s, x, y, w, h, fill=ARC_CARD, border=ARC_BORDER, border_w=1.0, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    c = s.shapes.add_shape(shape_type, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = border; c.line.width = Pt(border_w)
    c.shadow.inherit = False
    if radius:
        try:
            c.adjustments[0] = 0.06
        except Exception:
            pass
    return c


def pill(s, x, y, w, h, label, txt_color=ARC_BLUE, fill=None, border=ARC_BLUE):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill is None:
        p.fill.background()
    else:
        p.fill.solid(); p.fill.fore_color.rgb = fill
    p.line.color.rgb = border; p.line.width = Pt(1.0)
    p.shadow.inherit = False
    try:
        p.adjustments[0] = 0.5
    except Exception:
        pass
    tf = p.text_frame; tf.word_wrap = False
    tf.margin_top = 0; tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run(); r.text = label
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(r, 11, txt_color, bold=True)
    return p


def logo(s, cx, cy, d):
    """Reproduce logo.svg: blue circle + white plus."""
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
    circ.fill.solid(); circ.fill.fore_color.rgb = ARC_BLUE
    circ.line.fill.background(); circ.shadow.inherit = False
    t = int(d * 0.09)
    arm = int(d * 0.56)
    off = int((d - arm) / 2)
    thick = int(d * 0.145)
    # vertical bar
    v = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           cx + int(d/2 - thick/2), cy + off, thick, arm)
    # horizontal bar
    hb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            cx + off, cy + int(d/2 - thick/2), arm, thick)
    for bar in (v, hb):
        bar.fill.solid(); bar.fill.fore_color.rgb = WHITE
        bar.line.fill.background(); bar.shadow.inherit = False
    return circ


def page_tag(s, n, total, section):
    text(s, Inches(11.6), Inches(6.95), Inches(1.6), Inches(0.4),
         [[(f"{section}", 9, GRAY_500, False, True)]], align=PP_ALIGN.RIGHT)
    text(s, Inches(0.6), Inches(6.95), Inches(2.0), Inches(0.4),
         [[(f"{n:02d} / {total:02d}", 9, GRAY_500, False, True)]])


def accent_bar(s, x, y, h=Inches(0.42), w=Inches(0.06), color=ARC_BLUE):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background(); b.shadow.inherit = False
    return b


TOTAL = 10

# =============================================================================
# 1 — COVER
# =============================================================================
s = slide()
# subtle top accent line
accent_bar(s, Inches(0), Inches(0), h=Inches(0.09), w=EMU_W, color=ARC_BLUE)
logo(s, Inches(0.9), Inches(0.85), Inches(0.9))
text(s, Inches(1.95), Inches(0.92), Inches(5), Inches(0.8),
     [[("Cadence", 26, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)

pill(s, Inches(0.9), Inches(2.85), Inches(2.55), Inches(0.42),
     "LIVE ON ARC TESTNET", ARC_BLUE, fill=ARC_CARD, border=ARC_BLUE)

text(s, Inches(0.85), Inches(3.45), Inches(11.6), Inches(2.0),
     [[("Payroll that pays ", 52, WHITE, True), ("every second.", 52, ARC_BLUE, True)]],
     line_spacing=1.0)

text(s, Inches(0.9), Inches(5.05), Inches(9.8), Inches(1.2),
     [[("Continuous USDC payroll streaming on Arc. Employers deposit once, "
        "employees earn by the second and withdraw anytime in about 350ms.",
        17, GRAY_400)]], line_spacing=1.25)

text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
     [[("Hackathon checkpoint  ", 12, GRAY_500, True),
       ("|  Progress review  |  ", 12, GRAY_500),
       ("cadence-phi-ochre.vercel.app", 12, ARC_BLUE, False, True)]])

# =============================================================================
# 2 — THE PROBLEM
# =============================================================================
s = slide()
accent_bar(s, Inches(0.6), Inches(0.7))
text(s, Inches(0.8), Inches(0.62), Inches(11), Inches(0.6),
     [[("The problem", 30, WHITE, True)]])
text(s, Inches(0.82), Inches(1.35), Inches(11), Inches(0.6),
     [[("The monthly pay cycle is a leftover from paper banking.", 16, GRAY_400)]])

probs = [
    ("Workers wait", "People earn every day but get paid once a month. Money they have already earned sits locked in someone else's account."),
    ("Rails are slow", "Bank transfers take days and cross border payroll is worse. Fees stack up and cutoff times cause delays."),
    ("Trust is required", "Employees have to trust that the funds exist and will actually arrive. The balance is invisible until payday."),
]
cx = Inches(0.8)
cw = Inches(3.75)
gap = Inches(0.22)
for i, (h, d) in enumerate(probs):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    card(s, x, Inches(2.35), cw, Inches(3.4))
    text(s, Emu(int(x) + int(Inches(0.35))), Inches(2.7), Emu(int(cw) - int(Inches(0.7))), Inches(0.6),
         [[(f"0{i+1}", 13, ARC_BLUE, True, True)]])
    text(s, Emu(int(x) + int(Inches(0.35))), Inches(3.15), Emu(int(cw) - int(Inches(0.7))), Inches(0.7),
         [[(h, 20, WHITE, True)]])
    text(s, Emu(int(x) + int(Inches(0.35))), Inches(3.95), Emu(int(cw) - int(Inches(0.7))), Inches(1.7),
         [[(d, 14, GRAY_400)]], line_spacing=1.25)
page_tag(s, 2, TOTAL, "PROBLEM")

# =============================================================================
# 3 — THE SOLUTION
# =============================================================================
s = slide()
accent_bar(s, Inches(0.6), Inches(0.7))
text(s, Inches(0.8), Inches(0.62), Inches(11.5), Inches(0.6),
     [[("What Cadence does", 30, WHITE, True)]])
text(s, Inches(0.82), Inches(1.35), Inches(11.6), Inches(0.9),
     [[("Cadence flips payroll from a monthly event into a continuous flow. "
        "The moment an employer starts a stream, USDC begins moving to the "
        "employee every second, all on a smart contract.", 16, GRAY_400)]],
     line_spacing=1.25)

feats = [
    ("Settles in ~350ms", "Arc finality is sub second. Withdrawn funds arrive almost instantly."),
    ("Predictable fees", "Arc uses USDC for gas, so costs stay stable with no volatile gas token."),
    ("Withdraw anytime", "Employees pull what they have earned whenever they want. No payday."),
    ("Live runway", "Employers see exactly how long each stream can run before it empties."),
    ("Built in records", "Every stream carries an invoice or tax reference for clean bookkeeping."),
    ("Non custodial", "No company holds the money. Funds live in the contract until claimed."),
]
gx, gy = Inches(0.8), Inches(2.55)
cw, ch = Inches(3.75), Inches(1.9)
gapx, gapy = Inches(0.22), Inches(0.22)
for i, (h, d) in enumerate(feats):
    r, c = divmod(i, 3)
    x = Emu(int(gx) + c * (int(cw) + int(gapx)))
    y = Emu(int(gy) + r * (int(ch) + int(gapy)))
    card(s, x, y, cw, ch)
    text(s, Emu(int(x) + int(Inches(0.3))), Emu(int(y) + int(Inches(0.25))),
         Emu(int(cw) - int(Inches(0.6))), Inches(0.55),
         [[(h, 17, WHITE, True)]])
    text(s, Emu(int(x) + int(Inches(0.3))), Emu(int(y) + int(Inches(0.82))),
         Emu(int(cw) - int(Inches(0.6))), Inches(1.0),
         [[(d, 13, GRAY_400)]], line_spacing=1.2)
page_tag(s, 3, TOTAL, "SOLUTION")

# =============================================================================
# 4 — HOW IT WORKS
# =============================================================================
s = slide()
accent_bar(s, Inches(0.6), Inches(0.7))
text(s, Inches(0.8), Inches(0.62), Inches(11), Inches(0.6),
     [[("How it works", 30, WHITE, True)]])
text(s, Inches(0.82), Inches(1.35), Inches(11.6), Inches(0.6),
     [[("Five steps from deposit to withdrawal, fully on chain.", 16, GRAY_400)]])

steps = [
    ("01", "Approve and deposit", "The employer approves Cadence to spend USDC and deposits a lump sum."),
    ("02", "Create a stream", "Set the employee address, total amount, duration and an optional invoice reference."),
    ("03", "USDC flows live", "The contract streams per second. The employee dashboard ticks up in real time."),
    ("04", "Withdraw anytime", "The employee claims accrued earnings whenever they like. Funds land in ~350ms."),
    ("05", "Top up or cancel", "The employer adds funds or cancels. Unstreamed USDC returns to the employer."),
]
y = Inches(2.25)
for i, (n, h, d) in enumerate(steps):
    row_y = Emu(int(y) + i * int(Inches(0.92)))
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), row_y, Inches(0.6), Inches(0.6))
    num.fill.solid(); num.fill.fore_color.rgb = ARC_CARD
    num.line.color.rgb = ARC_BLUE; num.line.width = Pt(1.25); num.shadow.inherit = False
    ntf = num.text_frame; ntf.margin_top = 0; ntf.margin_bottom = 0
    nr = ntf.paragraphs[0].add_run(); nr.text = n
    ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(nr, 13, ARC_BLUE, bold=True, mono=True)
    text(s, Inches(1.75), Emu(int(row_y) - int(Inches(0.02))), Inches(3.4), Inches(0.6),
         [[(h, 18, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.2), Emu(int(row_y) - int(Inches(0.02))), Inches(7.4), Inches(0.6),
         [[(d, 14, GRAY_400)]], anchor=MSO_ANCHOR.MIDDLE)
page_tag(s, 4, TOTAL, "FLOW")

# =============================================================================
# 5 — PROGRESS (the checkpoint slide)
# =============================================================================
s = slide()
accent_bar(s, Inches(0.6), Inches(0.7))
text(s, Inches(0.8), Inches(0.62), Inches(11.5), Inches(0.6),
     [[("Where we are now", 30, WHITE, True)]])
text(s, Inches(0.82), Inches(1.35), Inches(11.6), Inches(0.6),
     [[("End to end product is live on Arc Testnet. Contract deployed, frontend shipped, "
        "core flows working.", 16, GRAY_400)]], line_spacing=1.2)

# metric strip
metrics = [("Deployed", "Contract live on Arc"), ("6 / 6", "Forge tests passing"),
           ("3", "App routes shipped"), ("4", "On chain actions working")]
mx, mw = Inches(0.8), Inches(2.92)
for i, (big, small) in enumerate(metrics):
    x = Emu(int(mx) + i * (int(mw) + int(Inches(0.18))))
    card(s, x, Inches(2.2), mw, Inches(1.15))
    text(s, Emu(int(x) + int(Inches(0.28))), Inches(2.35), Emu(int(mw) - int(Inches(0.5))), Inches(0.6),
         [[(big, 26, ARC_BLUE, True, True)]])
    text(s, Emu(int(x) + int(Inches(0.28))), Inches(2.92), Emu(int(mw) - int(Inches(0.5))), Inches(0.4),
         [[(small, 12, GRAY_400)]])

# Done column
done_x = Inches(0.8)
card(s, done_x, Inches(3.65), Inches(6.05), Inches(3.05))
text(s, Emu(int(done_x) + int(Inches(0.35))), Inches(3.85), Inches(5.4), Inches(0.5),
     [[("Shipped and working", 18, GREEN, True)]])
done = [
    "PayrollManager contract deployed to Arc Testnet",
    "Create, withdraw, top up and cancel all functioning",
    "Employer dashboard: balance, streams, create and manage",
    "Employee dashboard: live per second earnings ticker",
    "Wallet connect, USDC approval flow, live runway display",
    "Full Foundry test suite green (6 tests)",
]
text(s, Emu(int(done_x) + int(Inches(0.35))), Inches(4.35), Inches(5.4), Inches(2.3),
     [[("+  ", 13, GREEN, True), (d, 13.5, GRAY_400)] for d in done],
     line_spacing=1.15, space_after=6)

# In progress column
ip_x = Inches(7.1)
card(s, ip_x, Inches(3.65), Inches(5.45), Inches(3.05))
text(s, Emu(int(ip_x) + int(Inches(0.35))), Inches(3.85), Inches(4.8), Inches(0.5),
     [[("In progress and next", 18, ARC_BLUE, True)]])
nxt = [
    "Event indexing for withdrawal and payment history",
    "Contract audit pass and dust handling refinement",
    "Multi token support beyond USDC (EURC on Arc)",
    "Batch stream creation for teams",
    "Mainnet deployment and onboarding polish",
]
text(s, Emu(int(ip_x) + int(Inches(0.35))), Inches(4.35), Inches(4.8), Inches(2.3),
     [[(">  ", 13, ARC_BLUE, True), (d, 13.5, GRAY_400)] for d in nxt],
     line_spacing=1.15, space_after=6)
page_tag(s, 5, TOTAL, "PROGRESS")

# =============================================================================
# 6 — PRODUCT (what is built, described)
# =============================================================================
s = slide()
accent_bar(s, Inches(0.6), Inches(0.7))
text(s, Inches(0.8), Inches(0.62), Inches(11.5), Inches(0.6),
     [[("The product today", 30, WHITE, True)]])
text(s, Inches(0.82), Inches(1.35), Inches(11.6), Inches(0.6),
     [[("Two dashboards, one contract. Both are live in the deployed app.", 16, GRAY_400)]])

# Employer mock
ex = Inches(0.8)
ew = Inches(5.85)
card(s, ex, Inches(2.15), ew, Inches(4.55))
pill(s, Emu(int(ex) + int(Inches(0.35))), Inches(2.45), Inches(1.7), Inches(0.4),
     "EMPLOYER", ARC_BLUE, fill=ARC_DARK, border=ARC_BORDER)
text(s, Emu(int(ex) + int(Inches(0.35))), Inches(3.0), Inches(5.1), Inches(0.5),
     [[("Fund and manage payroll", 19, WHITE, True)]])
# inner mini balance card
mb = card(s, Emu(int(ex) + int(Inches(0.35))), Inches(3.65),
          Inches(5.15), Inches(1.05), fill=ARC_DARK)
text(s, Emu(int(ex) + int(Inches(0.6))), Inches(3.8), Inches(4), Inches(0.35),
     [[("USDC BALANCE", 10, GRAY_500, True)]])
text(s, Emu(int(ex) + int(Inches(0.6))), Inches(4.1), Inches(4), Inches(0.5),
     [[("$48,250.00", 24, WHITE, True, True)]])
for i, t in enumerate(["Create streams with rate and runway preview",
                       "Top up or cancel any active stream",
                       "See daily and monthly equivalents per stream"]):
    text(s, Emu(int(ex) + int(Inches(0.35))), Emu(int(Inches(4.95)) + i*int(Inches(0.5))),
         Inches(5.1), Inches(0.45),
         [[("+  ", 13, ARC_BLUE, True), (t, 13.5, GRAY_400)]])

# Employee mock
yx = Inches(6.95)
yw = Inches(5.6)
card(s, yx, Inches(2.15), yw, Inches(4.55))
pill(s, Emu(int(yx) + int(Inches(0.35))), Inches(2.45), Inches(1.7), Inches(0.4),
     "EMPLOYEE", GREEN, fill=ARC_DARK, border=ARC_BORDER)
text(s, Emu(int(yx) + int(Inches(0.35))), Inches(3.0), Inches(4.9), Inches(0.5),
     [[("Watch earnings in real time", 19, WHITE, True)]])
tk = card(s, Emu(int(yx) + int(Inches(0.35))), Inches(3.65),
          Inches(4.9), Inches(1.05), fill=ARC_DARK)
text(s, Emu(int(yx) + int(Inches(0.6))), Inches(3.8), Inches(4), Inches(0.35),
     [[("READY TO WITHDRAW", 10, GRAY_500, True)]])
text(s, Emu(int(yx) + int(Inches(0.6))), Inches(4.1), Inches(4.2), Inches(0.5),
     [[("$127", 24, WHITE, True, True), (".840512", 24, ARC_BLUE, True, True)]])
for i, t in enumerate(["Per second ticker interpolated at 60fps",
                       "One click withdraw of accrued USDC",
                       "Live runway and stream status"]):
    text(s, Emu(int(yx) + int(Inches(0.35))), Emu(int(Inches(4.95)) + i*int(Inches(0.5))),
         Inches(4.9), Inches(0.45),
         [[("+  ", 13, GREEN, True), (t, 13.5, GRAY_400)]])
page_tag(s, 6, TOTAL, "PRODUCT")

# =============================================================================
# 7 — ARCHITECTURE
# =============================================================================
s = slide()
accent_bar(s, Inches(0.6), Inches(0.7))
text(s, Inches(0.8), Inches(0.62), Inches(11), Inches(0.6),
     [[("Architecture", 30, WHITE, True)]])
text(s, Inches(0.82), Inches(1.35), Inches(11.6), Inches(0.6),
     [[("A thin, auditable contract with a real time frontend on top.", 16, GRAY_400)]])

layers = [
    ("FRONTEND", "Next.js 15  |  TypeScript  |  Tailwind",
     "Landing, employer and employee dashboards. Client side ticker interpolates on chain values so the UI stays smooth without hammering the RPC."),
    ("WEB3 LAYER", "wagmi v2  |  viem  |  injected wallet",
     "Typed read and write hooks with polling refetch. Approval and transaction flows handled with clear pending and error states."),
    ("CONTRACT", "Solidity 0.8  |  Foundry",
     "PayrollManager holds deposits and computes accrual per second. Reentrancy guarded. Deployed and verified on Arc Testnet."),
    ("NETWORK", "Arc Testnet  |  Chain 5042002",
     "USDC as the gas and payment token. Sub second finality makes streaming and instant withdrawal practical."),
]
y = Inches(2.2)
for i, (tag, stack, desc) in enumerate(layers):
    row_y = Emu(int(y) + i * int(Inches(1.12)))
    card(s, Inches(0.8), row_y, Inches(11.75), Inches(0.98))
    # tag chip
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Emu(int(row_y)+int(Inches(0.24))), Inches(1.85), Inches(0.5))
    chip.fill.solid(); chip.fill.fore_color.rgb = ARC_DARK
    chip.line.color.rgb = ARC_BLUE; chip.line.width = Pt(1.0); chip.shadow.inherit = False
    ctf = chip.text_frame; ctf.margin_top=0; ctf.margin_bottom=0
    cr = ctf.paragraphs[0].add_run(); cr.text = tag
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_font(cr, 11, ARC_BLUE, bold=True, mono=True)
    text(s, Inches(3.15), Emu(int(row_y)+int(Inches(0.14))), Inches(3.3), Inches(0.7),
         [[(stack, 12.5, WHITE, True, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(6.5), Emu(int(row_y)+int(Inches(0.1))), Inches(5.85), Inches(0.8),
         [[(desc, 12, GRAY_400)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
page_tag(s, 7, TOTAL, "ARCHITECTURE")

# =============================================================================
# 8 — CONTRACT
# =============================================================================
s = slide()
accent_bar(s, Inches(0.6), Inches(0.7))
text(s, Inches(0.8), Inches(0.62), Inches(11.5), Inches(0.6),
     [[("The contract", 30, WHITE, True)]])
text(s, Inches(0.82), Inches(1.35), Inches(11.6), Inches(0.6),
     [[("PayrollManager keeps the money logic small and legible.", 16, GRAY_400)]])

# left: functions
lx = Inches(0.8)
card(s, lx, Inches(2.2), Inches(6.0), Inches(4.5))
text(s, Emu(int(lx)+int(Inches(0.35))), Inches(2.4), Inches(5.4), Inches(0.5),
     [[("Core functions", 17, WHITE, True)]])
funcs = [
    ("createStream()", "employer opens a stream and funds it"),
    ("withdraw()", "employee claims accrued USDC"),
    ("topUp()", "employer adds funds, reactivates if drained"),
    ("cancelStream()", "pays accrued, refunds the rest"),
    ("accrued()", "view: earned since last claim"),
    ("runway()", "view: seconds of pay remaining"),
]
for i, (fn, d) in enumerate(funcs):
    yy = Emu(int(Inches(2.95)) + i*int(Inches(0.6)))
    text(s, Emu(int(lx)+int(Inches(0.35))), yy, Inches(2.65), Inches(0.5),
         [[(fn, 14, ARC_BLUE, True, True)]])
    text(s, Emu(int(lx)+int(Inches(3.0))), yy, Inches(3.15), Inches(0.5),
         [[(d, 12.5, GRAY_400)]])

# right: properties + address
rx = Inches(7.1)
card(s, rx, Inches(2.2), Inches(5.45), Inches(2.15))
text(s, Emu(int(rx)+int(Inches(0.35))), Inches(2.4), Inches(4.8), Inches(0.5),
     [[("Design choices", 17, WHITE, True)]])
props = [
    "Per second accrual, capped at the deposit",
    "Reentrancy guarded state changing calls",
    "Employer and employee scoped access checks",
    "Auto deactivates when a deposit runs dry",
]
text(s, Emu(int(rx)+int(Inches(0.35))), Inches(2.95), Inches(4.75), Inches(1.4),
     [[("- ", 12.5, ARC_BLUE, True), (p, 12.5, GRAY_400)] for p in props],
     line_spacing=1.2, space_after=4)

card(s, rx, Inches(4.55), Inches(5.45), Inches(2.15))
text(s, Emu(int(rx)+int(Inches(0.35))), Inches(4.75), Inches(4.8), Inches(0.4),
     [[("Deployed on Arc Testnet", 13, GRAY_500, True)]])
text(s, Emu(int(rx)+int(Inches(0.35))), Inches(5.15), Inches(4.85), Inches(0.4),
     [[("PayrollManager", 12, GRAY_400, True)]])
text(s, Emu(int(rx)+int(Inches(0.35))), Inches(5.5), Inches(4.85), Inches(0.4),
     [[("0x667Bc462AA0Bc3e2...A580D560", 12.5, ARC_BLUE, False, True)]])
text(s, Emu(int(rx)+int(Inches(0.35))), Inches(5.95), Inches(4.85), Inches(0.4),
     [[("USDC system token", 12, GRAY_400, True)]])
text(s, Emu(int(rx)+int(Inches(0.35))), Inches(6.3), Inches(4.85), Inches(0.4),
     [[("0x3600000000...00000000", 12.5, USDC_BLUE, False, True)]])
page_tag(s, 8, TOTAL, "CONTRACT")

# =============================================================================
# 9 — ROADMAP
# =============================================================================
s = slide()
accent_bar(s, Inches(0.6), Inches(0.7))
text(s, Inches(0.8), Inches(0.62), Inches(11), Inches(0.6),
     [[("Roadmap", 30, WHITE, True)]])
text(s, Inches(0.82), Inches(1.35), Inches(11.6), Inches(0.6),
     [[("From a working testnet product to production payroll infrastructure.", 16, GRAY_400)]])

phases = [
    ("NOW", "Checkpoint", GREEN,
     ["Contract live on Arc Testnet", "Both dashboards shipped", "Core flows working end to end", "Test suite passing"]),
    ("NEXT", "Hardening", ARC_BLUE,
     ["Payment history via event indexing", "Contract audit pass", "Dust and rounding refinement", "Onboarding and empty states"]),
    ("LATER", "Scale", GRAY_400,
     ["EURC and multi token streams", "Batch team payroll", "Mainnet deployment", "Employer analytics and exports"]),
]
cw = Inches(3.85)
for i, (tag, title, col, items) in enumerate(phases):
    x = Emu(int(Inches(0.8)) + i*(int(cw)+int(Inches(0.2))))
    card(s, x, Inches(2.3), cw, Inches(4.15))
    accent_bar(s, Emu(int(x)+int(Inches(0.35))), Inches(2.6), h=Inches(0.4), w=Inches(0.06), color=col)
    text(s, Emu(int(x)+int(Inches(0.55))), Inches(2.55), Inches(3), Inches(0.5),
         [[(tag, 12, col, True, True)]])
    text(s, Emu(int(x)+int(Inches(0.35))), Inches(3.1), Emu(int(cw)-int(Inches(0.7))), Inches(0.6),
         [[(title, 22, WHITE, True)]])
    text(s, Emu(int(x)+int(Inches(0.35))), Inches(3.8), Emu(int(cw)-int(Inches(0.7))), Inches(2.5),
         [[("- ", 13, col, True), (it, 13.5, GRAY_400)] for it in items],
         line_spacing=1.2, space_after=8)
page_tag(s, 9, TOTAL, "ROADMAP")

# =============================================================================
# 10 — CLOSE
# =============================================================================
s = slide()
accent_bar(s, Inches(0), Inches(0), h=Inches(0.09), w=EMU_W, color=ARC_BLUE)
logo(s, Inches(5.85), Inches(1.5), Inches(1.6))
text(s, Inches(0.5), Inches(3.25), Inches(12.33), Inches(0.9),
     [[("Cadence", 44, WHITE, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.6),
     [[("Real time payroll, live on Arc.", 20, ARC_BLUE, True)]], align=PP_ALIGN.CENTER)

card(s, Inches(3.15), Inches(5.15), Inches(7.0), Inches(1.4))
text(s, Inches(3.4), Inches(5.35), Inches(6.5), Inches(0.4),
     [[("Live app    ", 13, GRAY_500, True), ("cadence-phi-ochre.vercel.app", 13, ARC_BLUE, False, True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(3.4), Inches(5.75), Inches(6.5), Inches(0.4),
     [[("Code        ", 13, GRAY_500, True), ("github.com/dotmantissa/cadence", 13, WHITE, False, True)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(3.4), Inches(6.15), Inches(6.5), Inches(0.4),
     [[("Network     ", 13, GRAY_500, True), ("Arc Testnet, Chain 5042002", 13, GRAY_400, False, True)]],
     align=PP_ALIGN.CENTER)

# ---- save ---------------------------------------------------------------------
out = "deck/Cadence-Checkpoint-Deck.pptx"
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
